import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK_NAVY = RGBColor(15, 23, 42)
ROYAL_BLUE = RGBColor(30, 58, 138)
SLATE_GRAY = RGBColor(71, 85, 105)
WHITE = RGBColor(255, 255, 255)
LIGHT_BG = RGBColor(248, 250, 252)
ACCENT_GREEN = RGBColor(16, 185, 129)

def apply_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_slide_header(slide, title_text, is_dark=False):
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = "Segoe UI"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE if is_dark else DARK_NAVY
    
    # Thin divider line
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0.75), Inches(1.3), Inches(11.833), Inches(0.04)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_GREEN if is_dark else ROYAL_BLUE
    shape.line.fill.background()

# ==============================================================================
# SLIDE 1: TITLE SLIDE (Dark Navy Theme)
# ==============================================================================
slide_layout = prs.slide_layouts[6] # Blank
slide = prs.slides.add_slide(slide_layout)
apply_background(slide, DARK_NAVY)

title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.0))
tf = title_box.text_frame
tf.word_wrap = True

p1 = tf.paragraphs[0]
p1.text = "ONLINE SHOPPING MANAGEMENT SYSTEM"
p1.font.name = "Segoe UI"
p1.font.size = Pt(40)
p1.font.bold = True
p1.font.color.rgb = WHITE
p1.space_after = Pt(10)

p2 = tf.add_paragraph()
p2.text = "Python Tkinter Desktop Client & Oracle PL/SQL Integration"
p2.font.name = "Segoe UI"
p2.font.size = Pt(22)
p2.font.color.rgb = ACCENT_GREEN
p2.space_after = Pt(40)

p3 = tf.add_paragraph()
p3.text = "SIMAD UNIVERSITY - Faculty of Computing\nCourse: Oracle Database & PL/SQL Programming\nClient Application Platform: Standalone Python Tkinter GUI"
p3.font.name = "Segoe UI"
p3.font.size = Pt(14)
p3.font.color.rgb = RGBColor(148, 163, 184)

# ==============================================================================
# SLIDE 2: PROJECT OVERVIEW (Light Theme)
# ==============================================================================
slide = prs.slides.add_slide(slide_layout)
apply_background(slide, LIGHT_BG)
add_slide_header(slide, "1. Project Overview & Scope")

content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(5.0))
tf = content_box.text_frame
tf.word_wrap = True

points = [
    ("Objectives", "Develop a standalone desktop retail e-commerce management system (CARFOON shopping) that interfaces directly with Oracle XE."),
    ("Dual-Role Application", "Enables customers to browse items and purchase goods, and administrators to manage catalog items and track transactions in a single client interface."),
    ("Database-Centric Logic", "Ensures all critical calculations, constraints, inventory rules, and compliance logs reside and execute directly inside the database schema."),
    ("Primary Features", "Secure login panels, catalog search with image thumbnails, shopping cart sessions, PL/SQL bulk checkouts, receipts, and a Matplotlib dashboard.")
]

for title, desc in points:
    p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
    p.space_after = Pt(14)
    run_t = p.add_run()
    run_t.text = f"•  {title}: "
    run_t.font.bold = True
    run_t.font.size = Pt(16)
    run_t.font.name = "Segoe UI"
    run_t.font.color.rgb = ROYAL_BLUE
    
    run_d = p.add_run()
    run_d.text = desc
    run_d.font.size = Pt(15)
    run_d.font.name = "Segoe UI"
    run_d.font.color.rgb = SLATE_GRAY

# ==============================================================================
# SLIDE 3: DATABASE ARCHITECTURE (Light Theme)
# ==============================================================================
slide = prs.slides.add_slide(slide_layout)
apply_background(slide, LIGHT_BG)
add_slide_header(slide, "2. Oracle Database Architecture")

content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(5.0))
tf = content_box.text_frame
tf.word_wrap = True

points = [
    ("Oracle XE 1522 Instance", "The backend is powered by a local Oracle Database Express Edition (XE) instance listening on port 1522."),
    ("Dedicated User Schema", "All operations execute under the common user c##shop_user with explicit table privileges, ensuring strict isolation."),
    ("Thin Connection Model", "Leverages the python-oracledb Thin Mode driver. This creates direct socket communication from the Tkinter client to the database without requiring Oracle Instant Client binaries."),
    ("Transaction Safety", "Utilizes Oracle transaction controls, ensuring all DML operations execute atomically under connection-specific sessions.")
]

for title, desc in points:
    p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
    p.space_after = Pt(14)
    run_t = p.add_run()
    run_t.text = f"•  {title}: "
    run_t.font.bold = True
    run_t.font.size = Pt(16)
    run_t.font.name = "Segoe UI"
    run_t.font.color.rgb = ROYAL_BLUE
    
    run_d = p.add_run()
    run_d.text = desc
    run_d.font.size = Pt(15)
    run_d.font.name = "Segoe UI"
    run_d.font.color.rgb = SLATE_GRAY

# ==============================================================================
# SLIDE 4: DDL SCHEMA DESIGN (Light Theme)
# ==============================================================================
slide = prs.slides.add_slide(slide_layout)
apply_background(slide, LIGHT_BG)
add_slide_header(slide, "3. Relational Schema & Constraints")

content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(5.0))
tf = content_box.text_frame
tf.word_wrap = True

points = [
    ("Core Database Tables", "Enforces relations across 9 tables: USERS, PRODUCTS, CATEGORIES, CART, CART_ITEMS, ORDERS, ORDER_ITEMS, PAYMENTS, and ORDER_AUDIT_LOGS."),
    ("Referential Integrity", "Enforces foreign keys, primary keys, unique constraints (e.g., unique emails/product names), and CASCADE rules to prevent dangling references."),
    ("Normalized Design", "Follows Third Normal Form (3NF) to eliminate data anomalies and optimize insert/update speeds for catalog inventory and transaction records."),
    ("Audit Logs Tracking", "Keeps a detailed log history of database modifications, recording user context, dates, and actions for security compliance.")
]

for title, desc in points:
    p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
    p.space_after = Pt(14)
    run_t = p.add_run()
    run_t.text = f"•  {title}: "
    run_t.font.bold = True
    run_t.font.size = Pt(16)
    run_t.font.name = "Segoe UI"
    run_t.font.color.rgb = ROYAL_BLUE
    
    run_d = p.add_run()
    run_d.text = desc
    run_d.font.size = Pt(15)
    run_d.font.name = "Segoe UI"
    run_d.font.color.rgb = SLATE_GRAY

# ==============================================================================
# SLIDE 5: DATABASE TRIGGERS (Light Theme)
# ==============================================================================
slide = prs.slides.add_slide(slide_layout)
apply_background(slide, LIGHT_BG)
add_slide_header(slide, "4. Database Triggers & Automation")

content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(5.0))
tf = content_box.text_frame
tf.word_wrap = True

points = [
    ("Stock Adjustment (trg_update_stock)", "An AFTER INSERT trigger on ORDER_ITEMS. Instantly decrements product inventory levels in PRODUCTS when line items are checked out."),
    ("System Audit Logging (trg_system_audit)", "An AFTER INSERT OR UPDATE trigger on PRODUCTS that logs inventory movements directly into the ORDER_AUDIT_LOGS table."),
    ("Auto-Increment Generators", "BEFORE INSERT triggers on tables auto-fetch primary keys from Oracle sequence objects, simplifying insert queries."),
    ("Benefits", "Encapsulates workflow constraints inside the database catalog. This prevents the Tkinter client from having to manually calculate stock or log audit changes.")
]

for title, desc in points:
    p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
    p.space_after = Pt(14)
    run_t = p.add_run()
    run_t.text = f"•  {title}: "
    run_t.font.bold = True
    run_t.font.size = Pt(16)
    run_t.font.name = "Segoe UI"
    run_t.font.color.rgb = ROYAL_BLUE
    
    run_d = p.add_run()
    run_d.text = desc
    run_d.font.size = Pt(15)
    run_d.font.name = "Segoe UI"
    run_d.font.color.rgb = SLATE_GRAY

# ==============================================================================
# SLIDE 6: PL/SQL PACKAGES (Light Theme)
# ==============================================================================
slide = prs.slides.add_slide(slide_layout)
apply_background(slide, LIGHT_BG)
add_slide_header(slide, "5. Order Management PL/SQL Package")

content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(5.0))
tf = content_box.text_frame
tf.word_wrap = True

points = [
    ("Package-Based Encapsulation", "Collects database procedures and functions inside the PKG_ORDER_MANAGEMENT package, maintaining validation logic."),
    ("Procedure: add_item_to_cart", "Verifies stock levels and inserts new items into the CART_ITEMS table, raising a custom exception if stock is insufficient."),
    ("Procedure: process_bulk_checkout", "Converts cart items to an order, inserts order items, registers payments, and empties the cart atomically in a single session invocation."),
    ("Function: calculate_order_total", "Queries cart details, aggregates product prices, and calculates order totals, returning the numeric sum dynamically.")
]

for title, desc in points:
    p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
    p.space_after = Pt(14)
    run_t = p.add_run()
    run_t.text = f"•  {title}: "
    run_t.font.bold = True
    run_t.font.size = Pt(16)
    run_t.font.name = "Segoe UI"
    run_t.font.color.rgb = ROYAL_BLUE
    
    run_d = p.add_run()
    run_d.text = desc
    run_d.font.size = Pt(15)
    run_d.font.name = "Segoe UI"
    run_d.font.color.rgb = SLATE_GRAY

# ==============================================================================
# SLIDE 7: EXCEPTION HANDLING (Light Theme)
# ==============================================================================
slide = prs.slides.add_slide(slide_layout)
apply_background(slide, LIGHT_BG)
add_slide_header(slide, "6. Exception Handling & UI Validation")

content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(5.0))
tf = content_box.text_frame
tf.word_wrap = True

points = [
    ("Business Logic Safeguards", "Oracle raise_application_error triggers abort transaction actions if validation checks fail, maintaining integrity."),
    ("User-Defined Exception Tracing", "Package procedures define exceptions like e_invalid_quantity (PRAGMA EXCEPTION_INIT) to stop checkout if stocks are exceeded."),
    ("Graceful ORA-00001 Handling", "If an admin submits duplicate product or category names, the Tkinter app catches the database ORA-00001 constraint error and shows a readable warning."),
    ("Client-Side Translation", "The application translates raw database errors into user-friendly prompts, avoiding technical system crashes.")
]

for title, desc in points:
    p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
    p.space_after = Pt(14)
    run_t = p.add_run()
    run_t.text = f"•  {title}: "
    run_t.font.bold = True
    run_t.font.size = Pt(16)
    run_t.font.name = "Segoe UI"
    run_t.font.color.rgb = ROYAL_BLUE
    
    run_d = p.add_run()
    run_d.text = desc
    run_d.font.size = Pt(15)
    run_d.font.name = "Segoe UI"
    run_d.font.color.rgb = SLATE_GRAY

# ==============================================================================
# SLIDE 8: TKINTER UI IMPLEMENTATION (Light Theme)
# ==============================================================================
slide = prs.slides.add_slide(slide_layout)
apply_background(slide, LIGHT_BG)
add_slide_header(slide, "7. Desktop Client UI & Visual Polish")

content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(5.0))
tf = content_box.text_frame
tf.word_wrap = True

points = [
    ("Assign System Roles", "Only the primary administrator ('admin@gmail.com') is authorized to update roles. The primary 'admin@gmail.com' role itself is locked and non-editable."),
    ("User Details Editor", "Allows admins to update customer profiles (Full Name, Email Address, and Hashed Password) via sidebar entries."),
    ("Security Standards", "Enforces authentication checks. Default testing credentials are deleted so entry boxes load blank, requiring manual logins."),
    ("Interactive Subviews", "Includes multi-column catalog lists, an admin inventory editor, audit logs table grids, and popup transaction receipts.")
]

for title, desc in points:
    p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
    p.space_after = Pt(14)
    run_t = p.add_run()
    run_t.text = f"•  {title}: "
    run_t.font.bold = True
    run_t.font.size = Pt(16)
    run_t.font.name = "Segoe UI"
    run_t.font.color.rgb = ROYAL_BLUE
    
    run_d = p.add_run()
    run_d.text = desc
    run_d.font.size = Pt(15)
    run_d.font.name = "Segoe UI"
    run_d.font.color.rgb = SLATE_GRAY

# ==============================================================================
# SLIDE 9: IMAGE UPLOADS & VISUAL PREVIEWS (Light Theme)
# ==============================================================================
slide = prs.slides.add_slide(slide_layout)
apply_background(slide, LIGHT_BG)
add_slide_header(slide, "8. Product Images & Dynamic Previews")

content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(5.0))
tf = content_box.text_frame
tf.word_wrap = True

points = [
    ("Dynamic Copy & Db Sync", "Admins select images from file dialogs; the application copies them to local `/uploads` and saves the filename to the database."),
    ("In-Table Catalog Thumbnails", "The shopping tree catalog has 48px tall rows, displaying a 40x40 thumbnail of each product directly in the table list."),
    ("Layout Containers (tk.Frame)", "Wrapping the side preview labels in fixed 240x160 frames keeps aspect ratios intact and prevents window elements from shifting."),
    ("Click-to-Popup Full Resolution", "Hovering over the preview changes the cursor to a hand. Clicking the preview displays the full high-res original image in a centered window.")
]

for title, desc in points:
    p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
    p.space_after = Pt(14)
    run_t = p.add_run()
    run_t.text = f"•  {title}: "
    run_t.font.bold = True
    run_t.font.size = Pt(16)
    run_t.font.name = "Segoe UI"
    run_t.font.color.rgb = ROYAL_BLUE
    
    run_d = p.add_run()
    run_d.text = desc
    run_d.font.size = Pt(15)
    run_d.font.name = "Segoe UI"
    run_d.font.color.rgb = SLATE_GRAY

# ==============================================================================
# SLIDE 10: MATPLOTLIB ANALYTICS & CONCLUSION (Dark Navy Theme)
# ==============================================================================
slide = prs.slides.add_slide(slide_layout)
apply_background(slide, DARK_NAVY)
add_slide_header(slide, "9. Analytical Charts & Conclusion", is_dark=True)

content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(5.0))
tf = content_box.text_frame
tf.word_wrap = True

points = [
    ("Matplotlib UI Analytics", "Admins view a live business summary with real-time category sales distribution and order status donut charts."),
    ("Database-Centric System", "All constraints, triggers, and PL/SQL package subprograms operate directly inside Oracle XE, maintaining schema integrity."),
    ("Clean Tkinter Integrations", "Supports secure checkout, item addition, transaction processing, receipt generation, and audits in a clean desktop interface."),
    ("Reliability & Scalability", "The combination of thin-mode connection, encapsulated procedures, and automated triggers ensures a highly reliable e-commerce system.")
]

for title, desc in points:
    p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
    p.space_after = Pt(14)
    run_t = p.add_run()
    run_t.text = f"•  {title}: "
    run_t.font.bold = True
    run_t.font.size = Pt(16)
    run_t.font.name = "Segoe UI"
    run_t.font.color.rgb = ACCENT_GREEN
    
    run_d = p.add_run()
    run_d.text = desc
    run_d.font.size = Pt(15)
    run_d.font.name = "Segoe UI"
    run_d.font.color.rgb = RGBColor(203, 213, 225)

try:
    prs.save("Online_Shopping_Management_System_Presentation.pptx")
    print("Presentation generated successfully as 'Online_Shopping_Management_System_Presentation.pptx'!")
except PermissionError:
    prs.save("Online_Shopping_Management_System_Presentation_Updated.pptx")
    print("Presentation generated successfully as 'Online_Shopping_Management_System_Presentation_Updated.pptx' (original file was locked)!")
